import os
import json
import io
from typing import List, Dict, Any
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
import pdfplumber
from docx import Document
from pptx import Presentation
import pandas as pd
import logging

from config import GOOGLE_CLIENT_ID, GOOGLE_CLIENT_SECRET, GOOGLE_TOKEN_FILE

logger = logging.getLogger(__name__)

class GoogleDriveManager:
    def __init__(self):
        self.credentials = None
        self.service = None
        try:
            self._authenticate()
        except Exception as e:
            logger.error(f"Google Drive authentication failed: {e}")
            raise 

    def _authenticate(self):
        """Authenticate with Google Drive API using service account + domain-wide delegation"""
        from google.oauth2 import service_account

        SCOPES = ["https://www.googleapis.com/auth/drive"]
        SERVICE_ACCOUNT_FILE = "service_account.json"  # your key file path
        IMPERSONATED_USER = os.getenv("IMPERSONATED_USER", "investments@wyldvc.com")

        logger.info(f"Authenticating as delegated user: {IMPERSONATED_USER}")

        try:
            # Create credentials using service account and impersonated user
            creds = service_account.Credentials.from_service_account_file(
                SERVICE_ACCOUNT_FILE,
                scopes=SCOPES,
                subject=IMPERSONATED_USER
            )

            self.service = build("drive", "v3", credentials=creds)
            logger.info("Google Drive service initialized successfully (domain-wide delegation)")
        
        except Exception as e:
            logger.error(f"Error during authentication: {str(e)}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            raise



    async def download_dataroom_files(self) -> List[Dict[str, Any]]:
        """Download and parse all files (including subfolders) from the Google Drive dataroom."""
        if not self.service:
            raise ValueError("Google Drive service not initialized. Authentication may have failed.")

        try:
            drive_id = self._get_drive_id("Wyld VC")
            ir_folder_id = self._get_folder_id("IR", drive_id)
            dataroom_folder_id = self._get_folder_id("Wyld VC - Data Room", drive_id, parent_id=ir_folder_id)

            # 🔁 Recursively gather all files
            files = self._list_all_files_recursive(dataroom_folder_id, drive_id)
            logger.info(f"Found {len(files)} total files (including subfolders)")

            processed_files = []

            for file in files:
                try:
                    # Skip very large files
                    if file.get('size') and int(file['size']) > 50 * 1024 * 1024:
                        logger.warning(f"Skipping large file: {file['name']}")
                        continue

                    content = self._download_and_parse_file(file)
                    if content:
                        processed_files.append({
                            'id': file['id'],
                            'name': file['name'],
                            'content': content,
                            'mime_type': file['mimeType'],
                            'modified_time': file.get('modifiedTime')
                        })
                        logger.info(f"Processed: {file['name']}")

                except Exception as e:
                    logger.error(f"Error processing file {file.get('name', 'unknown')}: {str(e)}")
                    continue

            logger.info(f"Successfully processed {len(processed_files)} files (out of {len(files)})")
            return processed_files

        except Exception as e:
            logger.error(f"Error downloading files: {str(e)}")
            raise


    def _download_google_workspace_file(self, file_id: str, mime_type: str) -> str:
        """Download and convert Google Workspace files"""
        try:
            # Convert Google Docs/Sheets/Slides to text
            if mime_type == 'application/vnd.google-apps.document':
                # Export as plain text
                request = self.service.files().export_media(fileId=file_id, mimeType='text/plain')
            elif mime_type == 'application/vnd.google-apps.spreadsheet':
                # Export as CSV
                request = self.service.files().export_media(fileId=file_id, mimeType='text/csv')
            elif mime_type == 'application/vnd.google-apps.presentation':
                # Export as plain text
                request = self.service.files().export_media(fileId=file_id, mimeType='text/plain')
            else:
                logger.warning(f"Unsupported Google Workspace type: {mime_type}")
                return ""
            
            file_io = io.BytesIO()
            downloader = MediaIoBaseDownload(file_io, request)
            
            done = False
            while done is False:
                status, done = downloader.next_chunk()
            
            return file_io.getvalue().decode('utf-8')
            
        except Exception as e:
            logger.error(f"Error downloading Google Workspace file: {str(e)}")
            return ""

    def _parse_pdf(self, file_io: io.BytesIO) -> str:
        """Parse PDF content"""
        try:
            with pdfplumber.open(file_io) as pdf:
                text = ""
                for page in pdf.pages:
                    text += page.extract_text() or ""
                return text
        except Exception as e:
            logger.error(f"Error parsing PDF: {str(e)}")
            return ""

    def _parse_docx(self, file_io: io.BytesIO) -> str:
        """Parse DOCX content"""
        try:
            doc = Document(file_io)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error parsing DOCX: {str(e)}")
            return ""

    def _parse_pptx(self, file_io: io.BytesIO) -> str:
        """Parse PPTX content"""
        try:
            prs = Presentation(file_io)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error parsing PPTX: {str(e)}")
            return ""

    def _parse_excel_csv(self, file_io: io.BytesIO) -> str:
        """Parse Excel/CSV content"""
        try:
            df = pd.read_excel(file_io) if file_io.name.endswith('.xlsx') else pd.read_csv(file_io)
            return df.to_string()
        except Exception as e:
            logger.error(f"Error parsing Excel/CSV: {str(e)}")
            return ""


    def _get_drive_id(self, drive_name: str) -> str:
        response = self.service.drives().list().execute()
        for d in response.get('drives', []):
            if d['name'].lower() == drive_name.lower():
                return d['id']
        raise ValueError(f"Shared drive '{drive_name}' not found")


    def _get_folder_id(self, folder_name: str, drive_id: str, parent_id: str = None) -> str:
        """Get a folder ID by name, optionally within a parent folder."""
        query = f"mimeType='application/vnd.google-apps.folder' and name='{folder_name}'"
        if parent_id:
            query += f" and '{parent_id}' in parents"

        results = self.service.files().list(
            q=query,
            corpora="drive",
            driveId=drive_id,
            includeItemsFromAllDrives=True,
            supportsAllDrives=True,
            fields="files(id, name)"
        ).execute()

        folders = results.get('files', [])
        if not folders:
            raise ValueError(f"Folder '{folder_name}' not found (parent: {parent_id})")
        return folders[0]['id']

    #This prevents it just listing the files in the top level of the dataroom folder. It goes through all subfolders too.
    def _list_all_files_recursive(self, parent_id: str, drive_id: str) -> List[Dict[str, Any]]:
        """Recursively list all files under a folder in a shared drive."""
        all_files = []
        page_token = None

        while True:
            results = self.service.files().list(
                corpora="drive",
                driveId=drive_id,
                includeItemsFromAllDrives=True,
                supportsAllDrives=True,
                q=f"'{parent_id}' in parents and trashed=false",
                fields="nextPageToken, files(id, name, mimeType, size, modifiedTime)",
                pageToken=page_token
            ).execute()

            for f in results.get("files", []):
                if f["mimeType"] == "application/vnd.google-apps.folder":
                    all_files.extend(self._list_all_files_recursive(f["id"], drive_id))
                else:
                    all_files.append(f)

            page_token = results.get("nextPageToken")
            if not page_token:
                break

        return all_files

    def _download_and_parse_file(self, file: Dict[str, Any]) -> str:
        """Download and parse a single file based on its type"""
        file_id = file["id"]
        mime_type = file["mimeType"]

        try:
            # Handle Google Workspace files (Docs, Sheets, Slides)
            if mime_type.startswith("application/vnd.google-apps"):
                return self._download_google_workspace_file(file_id, mime_type)

            # Handle regular files
            request = self.service.files().get_media(fileId=file_id)
            file_io = io.BytesIO()
            downloader = MediaIoBaseDownload(file_io, request)

            done = False
            while not done:
                status, done = downloader.next_chunk()

            file_io.seek(0)

            # Parse based on file type
            if mime_type == "application/pdf":
                return self._parse_pdf(file_io)
            elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                return self._parse_docx(file_io)
            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                return self._parse_pptx(file_io)
            elif mime_type in [
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "text/csv",
            ]:
                return self._parse_excel_csv(file_io)
            elif mime_type.startswith("text/"):
                return file_io.read().decode("utf-8")
            else:
                logger.warning(f"Unsupported file type: {mime_type}")
                return ""

        except Exception as e:
            logger.error(f"Error downloading/parsing file {file.get('name', 'unknown')}: {str(e)}")
            return ""
